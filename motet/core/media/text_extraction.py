"""
Motet - Text Extraction Utilities

Copyright (c) 2024-2026 Motet Contributors
Licensed under the Functional Source License, Version 1.1, or a commercial license. See LICENSE.

Author: Matt Chisholm <matt@motet.dev>
Last Modified: 2026-08-24

Description:
    Utilities for extracting text from various file formats (PDF, DOCX, XLSX,
    PPTX, ODT, RTF).
    Note: PDF OCR is orchestrated by distributed commands:
    - `derive_pdf_page_images` stores page images as artifacts
    - `ocr_image_page` performs vision OCR per page
    This module stays "utilities-only" (no MotetContext, no distributed command composition).
    Used by the Derivation Pipeline.

Dependencies:
    - pdfplumber (PDF text extraction)
    - python-docx (DOCX)
    - python-pptx (PPTX)
    - odfpy (ODT)
    - striprtf (RTF)
    - openpyxl (XLSX)
    - io: Byte streams

Usage:
    text = extract_text_from_bytes(content, "application/pdf")
"""

import io
import structlog
import zipfile
import xml.etree.ElementTree as ET
from typing import Dict, Optional, Any

logger = structlog.get_logger(__name__)

def extract_text_from_bytes(content: bytes, content_type: str) -> str:
    """
    Extract text from raw file bytes based on content type.
    """
    if not content:
        return ""
        
    try:
        if content_type == "application/pdf":
            total_pages, layers = extract_pdf_text_layers(content)
            parts = []
            for page_num in range(1, total_pages + 1):
                text = (layers.get(page_num) or "").strip()
                if text:
                    parts.append(f"--- Page {page_num} (Text Layer) ---\n{text}")
            return "\n\n".join(parts)
        elif content_type in ["application/vnd.openxmlformats-officedocument.wordprocessingml.document", "application/msword"]:
            return _extract_docx(content)
        elif content_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
            return _extract_pptx(content)
        elif content_type == "application/vnd.oasis.opendocument.text":
            return _extract_odt(content)
        elif content_type in {"application/rtf", "text/rtf"}:
            return _extract_rtf(content)
        elif content_type in ["application/vnd.openxmlformats-officedocument.spreadsheetml.sheet", "application/vnd.ms-excel"]:
            return _extract_xlsx(content)
        elif content_type.startswith("text/") or content_type == "application/json" or content_type == "application/xml":
            return content.decode("utf-8", errors="ignore")
        else:
            logger.warning("unsupported_extraction_type", content_type=content_type)
            return ""
            
    except Exception as e:
        logger.error("extraction_failed", content_type=content_type, error=str(e))
        raise ValueError(f"Failed to extract text from {content_type}: {e}")

def extract_pdf_text_layers(pdf_content: bytes) -> tuple[int, Dict[int, str]]:
    """
    Extract per-page text from PDF text layers (no OCR).

    Returns:
        (total_pages, {page_num: text})
    """
    try:
        import pdfplumber
    except ImportError:
        logger.warning("pdfplumber_missing")
        return 0, {}

    try:
        with pdfplumber.open(io.BytesIO(pdf_content)) as pdf:
            total_pages = len(pdf.pages)
            layers: Dict[int, str] = {}
            for page_num, page in enumerate(pdf.pages, 1):
                text = page.extract_text()
                if text and text.strip():
                    layers[page_num] = text.strip()
            return total_pages, layers
    except Exception as e:
        logger.warning("pdf_text_layer_extraction_failed", error=str(e))
        return 0, {}


def combine_text_layer_and_ocr(page_num: int, text_layer: str, ocr_text: str) -> str:
    """
    Intelligently combine text layer and OCR results for a page.
    
    Strategy:
    - If text layer is substantial (>100 chars), use it as primary
    - If OCR finds significantly more text (>50% more), supplement with OCR
    - If text layer is sparse (<100 chars) but OCR found text, use OCR
    - If both exist, merge intelligently to avoid duplicates
    """
    has_text_layer = bool(text_layer and text_layer.strip())
    has_ocr = bool(ocr_text and ocr_text.strip())
    
    if not has_text_layer and not has_ocr:
        return ""
        
    # Case 1: Only text layer
    if has_text_layer and not has_ocr:
        return f"--- Page {page_num} (Text Layer) ---\n{text_layer}"
    
    # Case 2: Only OCR (scanned PDF)
    if not has_text_layer and has_ocr:
        return f"--- Page {page_num} (OCR) ---\n{ocr_text}"
    
    # Case 3: Both exist - combine intelligently
    text_layer_len = len(text_layer)
    ocr_len = len(ocr_text)
    
    # If text layer is substantial, use it as primary
    if text_layer_len > 100:
        # If OCR found significantly more (>50% more), supplement
        if ocr_len > text_layer_len * 1.5:
            logger.debug(
                "ocr_found_more_text_than_layer",
                page_num=page_num,
                text_layer_len=text_layer_len,
                ocr_len=ocr_len
            )
            # Combine both, with OCR as supplement
            return f"--- Page {page_num} (Text Layer + OCR Supplement) ---\n{text_layer}\n\n[Additional text from OCR:]\n{ocr_text}"
        else:
            # Text layer is primary and sufficient
            return f"--- Page {page_num} (Text Layer) ---\n{text_layer}"
    else:
        # Text layer is sparse, use OCR as primary
        logger.debug(
            "text_layer_sparse_using_ocr",
            page_num=page_num,
            text_layer_len=text_layer_len,
            ocr_len=ocr_len
        )
        if text_layer_len > 0:
            # Include sparse text layer for context
            return f"--- Page {page_num} (OCR + Sparse Text Layer) ---\n{ocr_text}\n\n[Text layer (sparse):]\n{text_layer}"
        else:
            return f"--- Page {page_num} (OCR) ---\n{ocr_text}"

def _extract_docx(content: bytes) -> str:
    try:
        import docx
    except ImportError:
        logger.warning("python_docx_missing")
        return ""
        
    doc = docx.Document(io.BytesIO(content))
    parts = [para.text for para in doc.paragraphs if para.text]
    for table_idx, table in enumerate(doc.tables, 1):
        parts.append(f"--- Table {table_idx} ---")
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            if any(cells):
                parts.append(" | ".join(cells))
    return "\n".join(parts)


def _extract_pptx(content: bytes) -> str:
    try:
        from pptx import Presentation
    except ImportError:
        logger.warning("python_pptx_missing")
        return _extract_pptx_ooxml(content)

    try:
        presentation = Presentation(io.BytesIO(content))
    except Exception as e:
        logger.warning("python_pptx_extraction_failed_using_ooxml_fallback", error=str(e))
        return _extract_pptx_ooxml(content)
    parts = []
    for slide_index, slide in enumerate(presentation.slides, 1):
        slide_parts = []
        for shape in slide.shapes:
            text = getattr(shape, "text", "")
            if text and text.strip():
                slide_parts.append(text.strip())
            if getattr(shape, "has_table", False):
                table = shape.table
                rows = []
                for row in table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    if any(cells):
                        rows.append(" | ".join(cells))
                if rows:
                    slide_parts.append("--- Table ---\n" + "\n".join(rows))
        notes = getattr(slide, "notes_slide", None)
        if notes is not None:
            note_text = getattr(getattr(notes, "notes_text_frame", None), "text", "")
            if note_text and note_text.strip():
                slide_parts.append(f"[Slide Notes]\n{note_text.strip()}")
        if slide_parts:
            parts.append(f"--- Page {slide_index} (Slide) ---\n" + "\n\n".join(slide_parts))
    extracted = "\n\n".join(parts)
    if extracted.strip():
        return extracted
    return _extract_pptx_ooxml(content)


def _extract_pptx_ooxml(content: bytes) -> str:
    """Fallback PPTX text extraction by reading slide and notes XML directly."""

    namespaces = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}

    def _sort_key(path: str) -> tuple[int, str]:
        stem = path.rsplit("/", 1)[-1].split(".", 1)[0]
        digits = "".join(ch for ch in stem if ch.isdigit())
        return (int(digits) if digits else 0, path)

    def _text_from_xml(raw_xml: bytes) -> str:
        try:
            root = ET.fromstring(raw_xml)
        except ET.ParseError:
            return ""
        texts = [
            (node.text or "").strip()
            for node in root.findall(".//a:t", namespaces)
            if (node.text or "").strip()
        ]
        return "\n".join(texts)

    parts: list[str] = []
    try:
        with zipfile.ZipFile(io.BytesIO(content)) as archive:
            slide_paths = sorted(
                (
                    name
                    for name in archive.namelist()
                    if name.startswith("ppt/slides/slide") and name.endswith(".xml")
                ),
                key=_sort_key,
            )
            for slide_index, name in enumerate(slide_paths, 1):
                text = _text_from_xml(archive.read(name))
                if text:
                    parts.append(f"--- Page {slide_index} (Slide) ---\n{text}")

            note_paths = sorted(
                (
                    name
                    for name in archive.namelist()
                    if name.startswith("ppt/notesSlides/notesSlide") and name.endswith(".xml")
                ),
                key=_sort_key,
            )
            for note_index, name in enumerate(note_paths, 1):
                text = _text_from_xml(archive.read(name))
                if text:
                    parts.append(f"--- Notes {note_index} ---\n{text}")
    except zipfile.BadZipFile as e:
        logger.warning("pptx_ooxml_extraction_bad_zip", error=str(e))
        return ""
    except Exception as e:
        logger.warning("pptx_ooxml_extraction_failed", error=str(e), exc_info=True)
        return ""

    return "\n\n".join(parts)


def _extract_odt(content: bytes) -> str:
    try:
        from odf import teletype
        from odf.opendocument import load
        from odf.table import Table, TableCell, TableRow
        from odf.text import P
    except ImportError:
        logger.warning("odfpy_missing")
        return ""

    doc = load(io.BytesIO(content))
    parts = []
    for paragraph in doc.getElementsByType(P):
        text = teletype.extractText(paragraph).strip()
        if text:
            parts.append(text)
    for table_idx, table in enumerate(doc.getElementsByType(Table), 1):
        parts.append(f"--- Table {table_idx} ---")
        for row in table.getElementsByType(TableRow):
            cells = []
            for cell in row.getElementsByType(TableCell):
                text = teletype.extractText(cell).strip()
                cells.append(text)
            if any(cells):
                parts.append(" | ".join(cells))
    return "\n".join(parts)


def _extract_rtf(content: bytes) -> str:
    try:
        from striprtf.striprtf import rtf_to_text
    except ImportError:
        logger.warning("striprtf_missing")
        return ""

    raw = content.decode("utf-8", errors="ignore")
    return rtf_to_text(raw)

def _extract_xlsx(content: bytes) -> str:
    try:
        import openpyxl
        import csv
    except ImportError:
        logger.warning("openpyxl_missing")
        return ""
        
    wb = openpyxl.load_workbook(io.BytesIO(content), read_only=True, data_only=True)
    parts = []
    
    for sheet_name in wb.sheetnames:
        parts.append(f"--- Sheet: {sheet_name} ---")
        sheet = wb[sheet_name]
        
        # Convert to CSV-like text
        rows = []
        for row in sheet.iter_rows(values_only=True):
            # Filter None
            clean_row = [str(cell) if cell is not None else "" for cell in row]
            if any(clean_row):
                rows.append(", ".join(clean_row))
        
        parts.append("\n".join(rows))
        
    return "\n\n".join(parts)


